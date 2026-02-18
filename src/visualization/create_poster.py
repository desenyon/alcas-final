"""
Generate Professional ISEF Poster for ALCAS
Clean 3-column layout with maximum visual impact
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Config
FIGURES_DIR = Path("results/figures")
ANALYSIS_DIR = Path("results/analysis")
POSTER_OUTPUT = Path("results/poster/ALCAS_ISEF_Poster.pptx")
POSTER_OUTPUT.parent.mkdir(parents=True, exist_ok=True)

# Professional ISEF Colors
PRIMARY = RGBColor(0, 51, 102)      # Navy blue
ACCENT = RGBColor(204, 102, 0)      # Burnt orange
LIGHT_BG = RGBColor(240, 248, 255)  # Alice blue
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

print("="*70)
print("CREATING PROFESSIONAL ISEF POSTER")
print("="*70)

# Create presentation (48" x 36")
prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

print("\nBuilding poster layout...")

# Helper functions
def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add colored rectangle"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, 
                bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add formatted text box"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color
    p.alignment = align
    
    return textbox

def add_section_box(slide, left, top, width, title, content, title_size=24, content_size=16):
    """Add section with title and content"""
    # Title background
    add_rectangle(slide, left, top, width, 0.5, PRIMARY)
    add_text_box(slide, left + 0.1, top + 0.05, width - 0.2, 0.4,
                title, font_size=title_size, bold=True, color=WHITE)
    
    # Content background
    add_rectangle(slide, left, top + 0.5, width, len(content.split('\n')) * 0.25 + 0.3, 
                 LIGHT_BG, PRIMARY)
    add_text_box(slide, left + 0.2, top + 0.7, width - 0.4, len(content.split('\n')) * 0.25,
                content, font_size=content_size, color=BLACK)

def add_image(slide, image_path, left, top, width=None, height=None):
    """Add image with error handling"""
    path = Path(image_path)
    if path.exists():
        try:
            if width and height:
                slide.shapes.add_picture(
                    str(path), Inches(left), Inches(top),
                    width=Inches(width), height=Inches(height)
                )
            elif width:
                slide.shapes.add_picture(
                    str(path), Inches(left), Inches(top), width=Inches(width)
                )
            else:
                slide.shapes.add_picture(str(path), Inches(left), Inches(top))
            return True
        except:
            return False
    return False

# ============================================================================
# HEADER SECTION
# ============================================================================
print("  Creating header...")

# Title banner
add_rectangle(slide, 0, 0, 48, 3, PRIMARY)

# Main title
add_text_box(slide, 1, 0.3, 46, 1.2,
            "ALCAS: Allosterically-Constrained Affinity Search",
            font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            font_name="Arial Black")

# Subtitle
add_text_box(slide, 1, 1.5, 46, 0.6,
            "Machine Learning for Rational Enzyme Design: Targeting Allosteric Sites to Engineer Superior PETase Variants",
            font_size=28, bold=False, color=RGBColor(200, 220, 255), align=PP_ALIGN.CENTER)

# Author info
add_text_box(slide, 1, 2.2, 46, 0.5,
            "Naitik Jariwala  •  ISEF 2025  •  Computational Biology & Bioinformatics  •  CMBIO047",
            font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================================
# LEFT COLUMN (16 inches wide)
# ============================================================================
LEFT_X = 0.5
COL_WIDTH = 15

print("  Building left column...")

# ABSTRACT
y_pos = 3.5
add_section_box(slide, LEFT_X, y_pos, COL_WIDTH, "ABSTRACT",
"""Plastic pollution threatens global ecosystems, with <9% of plastics recycled. PETase enzymes degrade PET plastic but are too slow for industrial use. Traditional enzyme engineering focuses on active-site mutations, yielding diminishing returns.

HYPOTHESIS: Allosteric mutations (distant from active site but dynamically coupled) outperform traditional active-site mutations.

APPROACH: Developed ALCAS combining graph neural networks (15,938 complexes, R²=0.448) with molecular dynamics to identify allosteric coupling. Designed 100 PETase variants under matched budgets.

RESULTS: Allosteric mutations outperformed active-site by 5.3% (p<0.001). Validated via ESMFold structures, statistical testing, and comprehensive ablations.

IMPACT: First ML-based allosteric enzyme design, opening new paradigms in protein engineering.""",
title_size=28, content_size=17)

# PROBLEM
y_pos = 10
add_section_box(slide, LEFT_X, y_pos, COL_WIDTH, "PROBLEM STATEMENT",
"""• 380 million tons plastic/year, <9% recycled
- PET degrades slowly (450+ years in environment)
- PETase can degrade PET but too inefficient
- Active-site engineering plateaued
- Need: Computational method for breakthrough improvement""",
title_size=24, content_size=18)

# INNOVATION
y_pos = 13.5
add_section_box(slide, LEFT_X, y_pos, COL_WIDTH, "INNOVATION & NOVELTY",
"""✓ First allosteric-constrained ML enzyme design
✓ Graph neural networks + molecular dynamics
✓ Matched-budget comparison (fair evaluation)
✓ 7 advanced extensions (transfer learning, Pareto, ESM-2)
✓ Comprehensive validation pipeline""",
title_size=24, content_size=18)

# FIGURE: 3D Mutation Regions
y_pos = 17
add_text_box(slide, LEFT_X, y_pos, COL_WIDTH, 0.4,
            "Figure 1: PETase Mutation Regions in 3D Space",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, FIGURES_DIR / "structures_3d/mutation_regions_3d.png",
         LEFT_X, y_pos + 0.5, width=COL_WIDTH)

# FIGURE: Distance Distribution
y_pos = 25
add_text_box(slide, LEFT_X, y_pos, COL_WIDTH, 0.4,
            "Figure 2: Allosteric Site Selection (Distance-Based Filtering)",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, FIGURES_DIR / "structures_3d/distance_distribution.png",
         LEFT_X, y_pos + 0.5, width=COL_WIDTH)

# ============================================================================
# MIDDLE COLUMN (16 inches wide)
# ============================================================================
MID_X = 16.5

print("  Building middle column...")

# METHODS
y_pos = 3.5
methods_content = """1. DATA PROCESSING
- PDBbind: 16,259 protein-ligand complexes
- Molecular graphs: 15,938 (atom-level features)
- Protein cluster splits (prevent leakage)

2. MODEL ARCHITECTURE
- Graph Isomorphism Networks (GIN)
- Virtual nodes + cross-attention
- 5 layers × 256 hidden dimensions
- 4-model ensemble (6.3M parameters)

3. ALLOSTERIC IDENTIFICATION
- Molecular dynamics (50ns, apo PETase)
- Network analysis (betweenness centrality)
- Dynamic coupling (correlation/MI)
- Criteria: >12Å distance + top coupling

4. VARIANT DESIGN
- 50 active + 50 allosteric variants
- Matched computational budget
- Ensemble scoring (μ ± σ)

5. VALIDATION
- ESMFold structure prediction
- Statistical testing (Mann-Whitney U)
- Ablation studies
- Baseline comparisons"""

add_section_box(slide, MID_X, y_pos, COL_WIDTH, "METHODS",
               methods_content, title_size=28, content_size=15)

# KEY EQUATIONS
y_pos = 15
add_text_box(slide, MID_X, y_pos, COL_WIDTH, 0.4,
            "Mathematical Framework",
            font_size=22, bold=True, color=PRIMARY)

add_rectangle(slide, MID_X, y_pos + 0.5, COL_WIDTH, 3.5, LIGHT_BG, PRIMARY)

equations = """Graph Neural Network Update:
h_i^(k+1) = MLP( h_i^(k) + Σ_j∈N(i) h_j^(k) )

Ensemble Prediction & Uncertainty:
μ = (1/M) Σ_m f_m(x)
σ² = (1/M) Σ_m (f_m(x) - μ)²

Pareto Dominance:
x ≻ y ⟺ ∀i: f_i(x) ≥ f_i(y) ∧ ∃j: f_j(x) > f_j(y)

Statistical Significance:
U = n₁n₂ + n₁(n₁+1)/2 - R₁
p-value from Mann-Whitney U distribution"""

add_text_box(slide, MID_X + 0.3, y_pos + 0.8, COL_WIDTH - 0.6, 3,
            equations, font_size=16, bold=True, color=BLACK, font_name="Courier New")

# FIGURE: Model Training
y_pos = 19.5
add_text_box(slide, MID_X, y_pos, COL_WIDTH, 0.4,
            "Figure 3: Model Training Performance (4-Model Ensemble)",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, FIGURES_DIR / "fig1_model_training.png",
         MID_X, y_pos + 0.5, width=COL_WIDTH)

# FIGURE: Pareto Optimization
y_pos = 27
add_text_box(slide, MID_X, y_pos, COL_WIDTH, 0.4,
            "Figure 4: Multi-Objective Pareto Optimization (52 Optimal Solutions)",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, ANALYSIS_DIR / "pareto/pareto_3d.png",
         MID_X, y_pos + 0.5, width=COL_WIDTH)

# ============================================================================
# RIGHT COLUMN (15 inches wide)
# ============================================================================
RIGHT_X = 32.5

print("  Building right column...")

# RESULTS
y_pos = 3.5
results_content = """PRIMARY FINDING
Allosteric > Active-site by 5.3%
p < 0.001 (Mann-Whitney U)

MODEL PERFORMANCE
- R² = 0.448 (test set)
- Pearson r = 0.673
- Calibrated uncertainty (r=0.113)

DESIGNED VARIANTS
- 100 total (50 active, 50 allosteric)
- 20 structures predicted (ESMFold)
- Allosteric: 7.09 ± 0.52 pKd
- Active-site: 6.73 ± 0.48 pKd
- Effect size: Cohen's d = 0.71

VALIDATION
- Ablations: All components essential
- vs Conservation: +40% (p<0.0001)
- Transfer learning: +6.3% R²"""

add_section_box(slide, RIGHT_X, y_pos, COL_WIDTH, "KEY RESULTS",
               results_content, title_size=28, content_size=17)

# FIGURE: ALCAS Results
y_pos = 11
add_text_box(slide, RIGHT_X, y_pos, COL_WIDTH, 0.4,
            "Figure 5: Allosteric vs Active-Site Comparison (Primary Finding)",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, FIGURES_DIR / "fig2_alcas_results.png",
         RIGHT_X, y_pos + 0.5, width=COL_WIDTH)

# FIGURE: Mutation Heatmap
y_pos = 18.5
add_text_box(slide, RIGHT_X, y_pos, COL_WIDTH, 0.4,
            "Figure 6: Complete Mutation Landscape (265 positions × 20 amino acids)",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, FIGURES_DIR / "heatmaps/full_mutation_landscape.png",
         RIGHT_X, y_pos + 0.5, width=COL_WIDTH)

# FIGURE: Summary Statistics
y_pos = 26
add_text_box(slide, RIGHT_X, y_pos, COL_WIDTH, 0.4,
            "Figure 7: Comprehensive Performance Summary",
            font_size=20, bold=True, color=PRIMARY)
add_image(slide, FIGURES_DIR / "fig3_summary_statistics.png",
         RIGHT_X, y_pos + 0.5, width=COL_WIDTH)

# ============================================================================
# BOTTOM SECTION (Full Width)
# ============================================================================
print("  Adding bottom sections...")

# ADVANCED EXTENSIONS
y_pos = 32.5
add_rectangle(slide, 0, y_pos, 48, 0.6, ACCENT)
add_text_box(slide, 1, y_pos + 0.05, 46, 0.5,
            "ADVANCED EXTENSIONS",
            font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

extensions = """• UNCERTAINTY-GUIDED DESIGN: Bayesian acquisition (UCB, Thompson sampling) for efficient exploration  • TRANSFER LEARNING: Fine-tuning improved R² by 6.3% on affinity extremes
- ESM-2 INTEGRATION: Tested protein language models (8,613 sequences) - structure alone sufficient  • PARETO OPTIMIZATION: 52 solutions balancing affinity/stability/solubility
- ABLATION STUDIES: GIN +82%, Virtual Nodes +138%, Cross-Attention +189%  • BASELINE COMPARISON: Outperforms evolutionary conservation by 40%"""

add_text_box(slide, 0.5, y_pos + 0.7, 47, 1, extensions,
            font_size=15, color=BLACK)

# CONCLUSIONS
y_pos = 34.2
add_rectangle(slide, 0, y_pos, 24, 0.5, PRIMARY)
add_text_box(slide, 0.5, y_pos + 0.05, 23, 0.4,
            "CONCLUSIONS & IMPACT",
            font_size=22, bold=True, color=WHITE)

conclusions = """• FIRST allosteric-constrained ML enzyme design
- Allosteric sites outperform active-site (+5.3%, p<0.001)
- Method generalizes (transfer learning validated)
- Accelerates plastic bioremediation
- Paradigm shift for rational protein engineering"""

add_text_box(slide, 0.5, y_pos + 0.6, 23, 1.2, conclusions,
            font_size=15, bold=True, color=BLACK)

# REFERENCES
add_rectangle(slide, 24, y_pos, 24, 0.5, PRIMARY)
add_text_box(slide, 24.5, y_pos + 0.05, 23, 0.4,
            "KEY REFERENCES",
            font_size=22, bold=True, color=WHITE)

refs = """[1] Austin et al. PNAS 2018 - PETase discovery  [2] Gilmer et al. ICML 2017 - Neural message passing
[3] Xu et al. ICLR 2019 - GIN theory  [4] Lin et al. bioRxiv 2022 - ESMFold  [5] Liu et al. - PDBbind database"""

add_text_box(slide, 24.5, y_pos + 0.6, 23, 1.2, refs,
            font_size=14, color=BLACK)

# Save
prs.save(str(POSTER_OUTPUT))

print(f"\n✓ Professional poster saved to: {POSTER_OUTPUT}")
print(f"  Size: 48\" × 36\"")
print(f"  Format: Editable PowerPoint")

print("\n" + "="*70)
print("POSTER COMPLETE")
print("="*70)
print("\nNext steps:")
print("1. Open in PowerPoint/Keynote")
print("2. Fine-tune figure positions")
print("3. Add school logo (top right)")
print("4. Export as PDF for printing")
print("5. Print at professional poster service (48\"×36\")")